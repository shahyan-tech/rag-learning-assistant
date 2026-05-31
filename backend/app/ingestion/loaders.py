from pathlib import Path
from typing import List

import nbformat # type: ignore
from pptx import Presentation # type: ignore
from pypdf import PdfReader # type: ignore
from langchain_core.documents import Document # type: ignore
from langchain_text_splitters import RecursiveCharacterTextSplitter # type: ignore


PROJECT_ROOT = Path(__file__).resolve().parents[3]
RAW_DATA_DIR = PROJECT_ROOT / "data" / "raw"

SUPPORTED_EXTENSIONS = {
    ".pdf",
    ".pptx",
    ".ipynb",
    ".txt",
    ".md",
}


def get_source_name(file_path: Path) -> str:
    """
    Return a clean source name relative to data/raw.

    Example:
    data/raw/ML Specialization/Course 1/Module 1/notes.pdf

    becomes:
    ML Specialization/Course 1/Module 1/notes.pdf
    """
    try:
        return file_path.relative_to(RAW_DATA_DIR).as_posix()
    except ValueError:
        return file_path.name


def load_pdf(file_path: Path) -> List[Document]:
    """Load text from a PDF file page by page."""
    documents = []
    reader = PdfReader(str(file_path))
    source_name = get_source_name(file_path)

    for page_number, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""

        if text.strip():
            documents.append(
                Document(
                    page_content=text.strip(),
                    metadata={
                        "source": source_name,
                        "file_name": file_path.name,
                        "file_path": str(file_path),
                        "file_type": "pdf",
                        "page": page_number,
                    },
                )
            )

    return documents


def load_pptx(file_path: Path) -> List[Document]:
    """Load text from a PowerPoint file slide by slide."""
    documents = []
    presentation = Presentation(str(file_path))
    source_name = get_source_name(file_path)

    for slide_number, slide in enumerate(presentation.slides, start=1):
        slide_text = []

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())

        combined_text = "\n".join(slide_text)

        if combined_text.strip():
            documents.append(
                Document(
                    page_content=combined_text.strip(),
                    metadata={
                        "source": source_name,
                        "file_name": file_path.name,
                        "file_path": str(file_path),
                        "file_type": "pptx",
                        "slide": slide_number,
                    },
                )
            )

    return documents


def load_ipynb(file_path: Path) -> List[Document]:
    """Load markdown and code cells from a Jupyter notebook."""
    documents = []
    notebook = nbformat.read(str(file_path), as_version=4)
    source_name = get_source_name(file_path)

    for cell_number, cell in enumerate(notebook.cells, start=1):
        cell_type = cell.get("cell_type", "")
        source = cell.get("source", "")

        if source.strip():
            documents.append(
                Document(
                    page_content=source.strip(),
                    metadata={
                        "source": source_name,
                        "file_name": file_path.name,
                        "file_path": str(file_path),
                        "file_type": "ipynb",
                        "cell": cell_number,
                        "cell_type": cell_type,
                    },
                )
            )

    return documents


def load_text_file(file_path: Path) -> List[Document]:
    """Load text from .txt or .md files."""
    text = file_path.read_text(encoding="utf-8", errors="ignore")
    source_name = get_source_name(file_path)

    if not text.strip():
        return []

    return [
        Document(
            page_content=text.strip(),
            metadata={
                "source": source_name,
                "file_name": file_path.name,
                "file_path": str(file_path),
                "file_type": file_path.suffix.lower().replace(".", ""),
            },
        )
    ]


def get_supported_files(raw_data_dir: Path = RAW_DATA_DIR) -> List[Path]:
    """
    Recursively find all supported files inside data/raw.

    This supports nested course folders:
    data/raw/Specialization/Course/Module/file.pdf
    """
    if not raw_data_dir.exists():
        raise FileNotFoundError(f"Raw data folder not found: {raw_data_dir}")

    files = []

    for file_path in raw_data_dir.rglob("*"):
        if file_path.is_dir() or file_path.name.startswith("."):
            continue

        if file_path.suffix.lower() in SUPPORTED_EXTENSIONS:
            files.append(file_path)

    return sorted(files)


def load_documents(raw_data_dir: Path = RAW_DATA_DIR) -> List[Document]:
    """Load all supported documents recursively from the raw data folder."""
    all_documents = []

    supported_files = get_supported_files(raw_data_dir)

    for file_path in supported_files:
        suffix = file_path.suffix.lower()

        if suffix == ".pdf":
            all_documents.extend(load_pdf(file_path))

        elif suffix == ".pptx":
            all_documents.extend(load_pptx(file_path))

        elif suffix == ".ipynb":
            all_documents.extend(load_ipynb(file_path))

        elif suffix in [".txt", ".md"]:
            all_documents.extend(load_text_file(file_path))

    return all_documents


def chunk_documents(documents: List[Document]) -> List[Document]:
    """Split documents into smaller chunks for RAG."""
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=1000,
        chunk_overlap=150,
        separators=["\n\n", "\n", ".", " ", ""],
    )

    chunks = text_splitter.split_documents(documents)
    return chunks